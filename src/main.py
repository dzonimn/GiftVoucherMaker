import os

import polars as pl
import qrcode
from pptx import Presentation

from misc import (
    delete_slide,
    duplicate_slide,
    replace_paragraph_text_retaining_initial_formatting,
)


def find_and_return_shape(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.rstrip() == text.rstrip():
            return shape
    raise Exception(
        f"Couldn't find a shape containing the text '{text}' in the template."
    )


if __name__ == "__main__":
    data = pl.scan_csv("./Hari Raya Vouchers.csv")
    data = data.select(["Voucher Code", "Voucher Link", "Name"]).collect()

    prs = Presentation("./Hari Raya Vouchers.pptx")
    for code, link, name in data.iter_rows():
        current_slide = duplicate_slide(prs, 0)
        code_shape = find_and_return_shape(current_slide, "Voucher Code")
        link_shape = find_and_return_shape(current_slide, "Voucher Link")
        name_shape = find_and_return_shape(current_slide, "NAME")
        qrcode_shape = find_and_return_shape(current_slide, "QR Code\n\n")

        replace_paragraph_text_retaining_initial_formatting(
            code_shape.text_frame.paragraphs[0], code
        )
        replace_paragraph_text_retaining_initial_formatting(
            link_shape.text_frame.paragraphs[0], link
        )
        replace_paragraph_text_retaining_initial_formatting(
            name_shape.text_frame.paragraphs[0], name
        )
        qrcode.make(code).save("qrcode.png")
        current_slide.shapes.add_picture(
            "./qrcode.png", qrcode_shape.left, qrcode_shape.top, qrcode_shape.height
        )
        os.remove("qrcode.png")

    delete_slide(prs, prs.slides[0])
    prs.save("output.pptx")
