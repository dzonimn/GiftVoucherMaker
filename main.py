from pptx import Presentation
import qrcode
import csv
import os

from misc import (
    delete_slide,
    duplicate_slide,
    replace_paragraph_text_retaining_initial_formatting,
)


def find_and_return_shape(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.rstrip() == text.rstrip():
            return shape
    raise Exception(
        f"Couldn't find a shape containing the text '{text}' in the template."
    )


if __name__ == "__main__":
    voucher_codes = []
    voucher_links = []
    with open(
        "Hari Raya Vouchers.csv",
    ) as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            # Skip the first line
            if reader.line_num == 1:
                continue
            # Otherwise, store the code and link
            voucher_code, voucher_link = row
            voucher_codes.append(voucher_code)
            voucher_links.append(voucher_link)

    prs = Presentation("./template_hariraya.pptx")
    for code, link in zip(voucher_codes, voucher_links):
        current_slide = duplicate_slide(prs, 0)
        code_shape = find_and_return_shape(current_slide, "Voucher Code")
        link_shape = find_and_return_shape(current_slide, "Voucher Link")
        qrcode_shape = find_and_return_shape(current_slide, "QR Code\n\n")

        replace_paragraph_text_retaining_initial_formatting(
            code_shape.text_frame.paragraphs[0], code
        )
        replace_paragraph_text_retaining_initial_formatting(
            link_shape.text_frame.paragraphs[0], link
        )
        qrcode.make(code).save("qrcode.png")
        current_slide.shapes.add_picture(
            "./qrcode.png", qrcode_shape.left, qrcode_shape.top, qrcode_shape.height
        )
        os.remove("qrcode.png")

    delete_slide(prs, prs.slides[0])
    prs.save("output.pptx")
